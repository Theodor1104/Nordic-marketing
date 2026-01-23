from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# Farver
PRIMARY = RGBColor(26, 61, 46)
PRIMARY_LIGHT = RGBColor(45, 90, 69)
ACCENT = RGBColor(61, 139, 110)
ACCENT_LIGHT = RGBColor(91, 168, 138)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 249)
GRAY = RGBColor(107, 125, 115)
DARK_GRAY = RGBColor(74, 90, 81)

# Service colors
META_BLUE = RGBColor(8, 102, 255)
META_PURPLE = RGBColor(160, 51, 255)
GOOGLE_BLUE = RGBColor(66, 133, 244)
GOOGLE_RED = RGBColor(234, 67, 53)
WEB_PURPLE = RGBColor(102, 126, 234)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_gradient_bg(slide, color1, color2):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color1
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color2
    accent.fill.fore_color.brightness = 0.1
    accent.line.fill.background()

    accent2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(5), Inches(5))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = color2
    accent2.fill.fore_color.brightness = 0.1
    accent2.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return box

def add_meta_icon(slide, x, y, size, bg_color):
    """Facebook 'f' ikon"""
    # Baggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.adjustments[0] = 0.2

    # F bogstav
    add_text_box(slide, x, y + size*0.15, size, size*0.8, "f", int(size*35), WHITE, True, PP_ALIGN.CENTER)

def add_google_icon(slide, x, y, size, bg_color):
    """Forstørrelsesglas med plus"""
    # Baggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.adjustments[0] = 0.2

    # Cirkel (lup)
    circle_size = size * 0.45
    circle_x = x + size * 0.2
    circle_y = y + size * 0.2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(circle_x), Inches(circle_y), Inches(circle_size), Inches(circle_size))
    circle.fill.background()
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(3)

    # Håndtag
    handle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + size*0.55), Inches(y + size*0.6), Inches(size*0.25), Inches(size*0.08))
    handle.fill.solid()
    handle.fill.fore_color.rgb = WHITE
    handle.line.fill.background()
    handle.rotation = 45

    # Plus (vandret)
    plus_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(circle_x + circle_size*0.25), Inches(circle_y + circle_size*0.42), Inches(circle_size*0.5), Inches(circle_size*0.16))
    plus_h.fill.solid()
    plus_h.fill.fore_color.rgb = WHITE
    plus_h.line.fill.background()

    # Plus (lodret)
    plus_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(circle_x + circle_size*0.42), Inches(circle_y + circle_size*0.25), Inches(circle_size*0.16), Inches(circle_size*0.5))
    plus_v.fill.solid()
    plus_v.fill.fore_color.rgb = WHITE
    plus_v.line.fill.background()

def add_seo_icon(slide, x, y, size, bg_color):
    """Søjlediagram"""
    # Baggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.adjustments[0] = 0.2

    # Søjler
    bar_width = size * 0.15
    gap = size * 0.08
    start_x = x + size * 0.2

    heights = [0.3, 0.5, 0.7]  # Stigende højder
    for i, h in enumerate(heights):
        bar_x = start_x + i * (bar_width + gap)
        bar_height = size * h
        bar_y = y + size * 0.75 - bar_height
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bar_x), Inches(bar_y), Inches(bar_width), Inches(bar_height))
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE
        bar.line.fill.background()

    # Bundlinje
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + size*0.15), Inches(y + size*0.75), Inches(size*0.7), Inches(size*0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

def add_web_icon(slide, x, y, size, bg_color):
    """Browser/skærm ikon"""
    # Baggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.adjustments[0] = 0.2

    # Skærm
    screen_w = size * 0.7
    screen_h = size * 0.45
    screen_x = x + size * 0.15
    screen_y = y + size * 0.2
    screen = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(screen_x), Inches(screen_y), Inches(screen_w), Inches(screen_h))
    screen.fill.background()
    screen.line.color.rgb = WHITE
    screen.line.width = Pt(2.5)
    screen.adjustments[0] = 0.1

    # Browser bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(screen_x), Inches(screen_y), Inches(screen_w), Inches(size*0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()

    # Prikker i browser bar
    for i in range(2):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(screen_x + 0.06 + i*0.08), Inches(screen_y + 0.03), Inches(0.05), Inches(0.05))
        dot.fill.solid()
        dot.fill.fore_color.rgb = bg_color
        dot.line.fill.background()

    # Stand
    stand = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + size*0.45), Inches(y + size*0.65), Inches(size*0.1), Inches(size*0.12))
    stand.fill.solid()
    stand.fill.fore_color.rgb = WHITE
    stand.line.fill.background()

    # Base
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + size*0.3), Inches(y + size*0.77), Inches(size*0.4), Inches(size*0.05))
    base.fill.solid()
    base.fill.fore_color.rgb = WHITE
    base.line.fill.background()

def add_shield_icon(slide, x, y, size, bg_color):
    """Skjold med flueben"""
    # Baggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.adjustments[0] = 0.2

    # Skjold (brug chevron som approksimation)
    shield = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(x + size*0.25), Inches(y + size*0.15), Inches(size*0.5), Inches(size*0.6))
    shield.fill.background()
    shield.line.color.rgb = WHITE
    shield.line.width = Pt(2.5)
    shield.rotation = 180

    # Flueben
    add_text_box(slide, x, y + size*0.2, size, size*0.6, "✓", int(size*28), WHITE, True, PP_ALIGN.CENTER)

def add_grid_icon(slide, x, y, size, bg_color):
    """Fire firkanter (komplet løsning)"""
    # Baggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.adjustments[0] = 0.2

    # Fire firkanter
    box_size = size * 0.25
    gap = size * 0.1
    start = size * 0.2

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for col, row in positions:
        box_x = x + start + col * (box_size + gap)
        box_y = y + start + row * (box_size + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box_x), Inches(box_y), Inches(box_size), Inches(box_size))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.fill.background()
        box.adjustments[0] = 0.15

def add_phone_icon(slide, x, y, size, bg_color):
    """Smartphone ikon"""
    # Baggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.adjustments[0] = 0.2

    # Telefon
    phone_w = size * 0.35
    phone_h = size * 0.6
    phone_x = x + size * 0.325
    phone_y = y + size * 0.2
    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(phone_x), Inches(phone_y), Inches(phone_w), Inches(phone_h))
    phone.fill.background()
    phone.line.color.rgb = WHITE
    phone.line.width = Pt(2.5)
    phone.adjustments[0] = 0.15

    # Hjem-knap
    btn = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(phone_x + phone_w*0.35), Inches(phone_y + phone_h*0.85), Inches(phone_w*0.3), Inches(phone_w*0.3))
    btn.fill.solid()
    btn.fill.fore_color.rgb = WHITE
    btn.line.fill.background()

def add_heart_icon(slide, x, y, size, bg_color):
    """Hjerte ikon"""
    # Baggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.adjustments[0] = 0.2

    # Hjerte med tekst
    add_text_box(slide, x, y + size*0.18, size, size*0.7, "♥", int(size*32), WHITE, False, PP_ALIGN.CENTER)

# ============ SLIDE 1: FORSIDE ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, PRIMARY, ACCENT)

logo_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.666), Inches(1.5), Inches(2), Inches(2))
logo_circle.fill.solid()
logo_circle.fill.fore_color.rgb = WHITE
logo_circle.fill.fore_color.brightness = 0.9
logo_circle.line.fill.background()

add_text_box(slide, 5.666, 1.9, 2, 1.2, "N", 72, PRIMARY, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 4, 12.333, 1.2, "Nordic Marketing", 60, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 5.2, 12.333, 0.8, "Digital marketing der skaber resultater", 28, ACCENT_LIGHT, False, PP_ALIGN.CENTER)

# ============ SLIDE 2: HVEM ER VI ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Hvem er vi?", 44, WHITE, True)

for i, (x, name) in enumerate([(2.5, "Grundlægger 1"), (8, "Grundlægger 2")]):
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.3), Inches(2.5), Inches(2.5))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = ACCENT
    avatar.line.fill.background()

    # Person ikon i avatar
    add_text_box(slide, x, 2.8, 2.5, 1.5, "👤", 50, WHITE, False, PP_ALIGN.CENTER)

    add_text_box(slide, x - 0.5, 5, 3.5, 0.6, name, 24, PRIMARY, True, PP_ALIGN.CENTER)
    add_text_box(slide, x - 0.5, 5.5, 3.5, 0.5, "Niels Brock", 18, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 6.3, 10.333, 0.8, "To unge iværksættere med passion for digital marketing - vokset op i den digitale verden", 20, GRAY, False, PP_ALIGN.CENTER)

# ============ SLIDE 3: PROBLEMET ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Problemet", 44, WHITE, True)

for i, (x, number, label) in enumerate([
    (1.5, "73%", "af små virksomheder\nkæmper med online\nmarkedsføring"),
    (5.5, "10K+", "kroner koster\ntraditionelle bureauer\nom måneden"),
    (9.5, "0", "gennemsigtighed\ni branchen")
]):
    stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(3), Inches(3.5))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = OFF_WHITE
    stat_box.line.color.rgb = RGBColor(230, 235, 232)
    stat_box.adjustments[0] = 0.1

    add_text_box(slide, x, 2.6, 3, 1, number, 48, ACCENT, True, PP_ALIGN.CENTER)

    label_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.8), Inches(2.6), Inches(1.5))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 4: LØSNINGEN ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Vores løsning", 44, WHITE, True)

solutions = [
    ("Gennemsigtig marketing", "til fair priser"),
    ("No cure, no pay", "du betaler for resultater"),
    ("Moderne værktøjer", "og intelligent brug af AI"),
    ("Personlig service", "hvor du er prioriteten")
]

for i, (title, subtitle) in enumerate(solutions):
    y = 2.2 + (i * 1.2)

    icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(0.6), Inches(0.6))
    icon.fill.solid()
    icon.fill.fore_color.rgb = ACCENT
    icon.line.fill.background()
    icon.adjustments[0] = 0.25

    add_text_box(slide, 1, y + 0.05, 0.6, 0.5, "✓", 24, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, 2, y, 5, 0.6, title, 24, PRIMARY, True)
    add_text_box(slide, 2, y + 0.45, 5, 0.5, subtitle, 18, ACCENT, False)

# ============ SLIDE 5: SERVICES ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Hvad vi tilbyder", 44, WHITE, True)

services = [
    (0.8, 2.2, "Meta Ads", "Facebook & Instagram annoncer", META_BLUE, "meta"),
    (4.05, 2.2, "Google Ads", "Søgemaskine annoncering", GOOGLE_BLUE, "google"),
    (7.3, 2.2, "SEO", "Søgemaskineoptimering", PRIMARY, "seo"),
    (10.55, 2.2, "Hjemmesider", "Moderne webdesign", WEB_PURPLE, "web")
]

for x, y, title, desc, color, icon_type in services:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.75), Inches(4))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(230, 235, 232)
    card.adjustments[0] = 0.08

    # Tilføj det rigtige ikon
    icon_x = x + 0.8
    icon_y = y + 0.5
    icon_size = 1.15

    if icon_type == "meta":
        add_meta_icon(slide, icon_x, icon_y, icon_size, color)
    elif icon_type == "google":
        add_google_icon(slide, icon_x, icon_y, icon_size, color)
    elif icon_type == "seo":
        add_seo_icon(slide, icon_x, icon_y, icon_size, color)
    elif icon_type == "web":
        add_web_icon(slide, icon_x, icon_y, icon_size, color)

    add_text_box(slide, x, y + 2, 2.75, 0.6, title, 22, PRIMARY, True, PP_ALIGN.CENTER)

    desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.6), Inches(2.35), Inches(1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 6: HVORFOR OS ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Hvorfor vælge os?", 44, WHITE, True)

usps = [
    (0.8, 2.2, "No cure, no pay", "Du betaler kun for resultater", "shield"),
    (7, 2.2, "Komplet pakke", "Marketing OG hjemmeside", "grid"),
    (0.8, 4.5, "Digital natives", "Vi forstår online trends", "phone"),
    (7, 4.5, "Personlig service", "Du er vores førsteprioritet", "heart")
]

for x, y, title, desc, icon_type in usps:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.5), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(230, 235, 232)
    card.adjustments[0] = 0.06

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.15), Inches(0.12), Inches(1.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()

    # Tilføj det rigtige ikon
    icon_x = x + 0.4
    icon_y = y + 0.4
    icon_size = 1.0

    if icon_type == "shield":
        add_shield_icon(slide, icon_x, icon_y, icon_size, PRIMARY)
    elif icon_type == "grid":
        add_grid_icon(slide, icon_x, icon_y, icon_size, PRIMARY)
    elif icon_type == "phone":
        add_phone_icon(slide, icon_x, icon_y, icon_size, PRIMARY)
    elif icon_type == "heart":
        add_heart_icon(slide, icon_x, icon_y, icon_size, PRIMARY)

    add_text_box(slide, x + 1.6, y + 0.4, 3.5, 0.5, title, 22, PRIMARY, True)
    add_text_box(slide, x + 1.6, y + 0.95, 3.5, 0.5, desc, 16, GRAY, False)

# ============ SLIDE 7: PROCESSEN ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Sådan arbejder vi", 44, WHITE, True)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.3), Inches(9.5), Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

steps = [
    (1.3, "1", "Gratis\nkonsultation"),
    (4.3, "2", "Strategi &\nforslag"),
    (7.3, "3", "Justering"),
    (10.3, "4", "Resultater")
]

for x, num, label in steps:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.5), Inches(1.7), Inches(1.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()

    add_text_box(slide, x, 2.85, 1.7, 1, num, 40, WHITE, True, PP_ALIGN.CENTER)

    label_box = slide.shapes.add_textbox(Inches(x - 0.3), Inches(4.5), Inches(2.3), Inches(1.2))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 8: CTA ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, PRIMARY, ACCENT)

add_text_box(slide, 0.5, 2.2, 12.333, 1, "Klar til at vokse online?", 52, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 3.5, 12.333, 0.8, "Book en gratis og uforpligtende konsultation", 26, ACCENT_LIGHT, False, PP_ALIGN.CENTER)

contacts = [
    (3, "kontakt@nordicmarketing.dk"),
    (7.5, "nordicmarketing.dk")
]

for x, text in contacts:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5), Inches(3.5), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.fill.fore_color.brightness = 0.9
    box.line.fill.background()
    box.adjustments[0] = 0.3

    add_text_box(slide, x, 5.2, 3.5, 0.5, text, 16, PRIMARY, False, PP_ALIGN.CENTER)

# Gem
output_path = '/Users/theodorhauch/nordic-marketing/Nordic_Marketing_Pitch.pptx'
prs.save(output_path)
print(f"PowerPoint gemt: {output_path}")
