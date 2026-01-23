from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Farver
PRIMARY = RGBColor(26, 61, 46)
PRIMARY_LIGHT = RGBColor(45, 90, 69)
ACCENT = RGBColor(61, 139, 110)
ACCENT_LIGHT = RGBColor(91, 168, 138)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 249)
GRAY = RGBColor(107, 125, 115)
DARK_GRAY = RGBColor(74, 90, 81)

# Meta colors
META_BLUE = RGBColor(8, 102, 255)
META_PURPLE = RGBColor(160, 51, 255)
GOOGLE_BLUE = RGBColor(66, 133, 244)
GOOGLE_RED = RGBColor(234, 67, 53)
WEB_PURPLE = RGBColor(102, 126, 234)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_gradient_bg(slide, color1, color2):
    """Tilføj baggrund med gradient-lignende effekt"""
    # Hovedbaggrund
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color1
    bg.line.fill.background()

    # Accent shape i hjørnet for visual interest
    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color2
    accent.fill.fore_color.brightness = 0.1
    accent.line.fill.background()

    # Nederste accent
    accent2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(5), Inches(5))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = color2
    accent2.fill.fore_color.brightness = 0.1
    accent2.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT):
    """Hjælpefunktion til tekstboks"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return box

def add_icon_box(slide, left, top, size, color):
    """Tilføj farvet ikon-boks"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(size))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    # Afrunding
    box.adjustments[0] = 0.2
    return box

# ============ SLIDE 1: FORSIDE ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, PRIMARY, ACCENT)

# Logo cirkel
logo_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.666), Inches(1.5), Inches(2), Inches(2))
logo_circle.fill.solid()
logo_circle.fill.fore_color.rgb = WHITE
logo_circle.fill.fore_color.brightness = 0.9
logo_circle.line.fill.background()

# Logo tekst
add_text_box(slide, 5.666, 1.9, 2, 1.2, "N", 72, PRIMARY, True, PP_ALIGN.CENTER)

# Titel
add_text_box(slide, 0.5, 4, 12.333, 1.2, "Nordic Marketing", 60, WHITE, True, PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, 0.5, 5.2, 12.333, 0.8, "Digital marketing der skaber resultater", 28, ACCENT_LIGHT, False, PP_ALIGN.CENTER)

# ============ SLIDE 2: HVEM ER VI ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Hvem er vi?", 44, WHITE, True)

# Team member boxes
for i, (x, name) in enumerate([(2.5, "Grundlægger 1"), (8, "Grundlægger 2")]):
    # Avatar cirkel
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.3), Inches(2.5), Inches(2.5))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = ACCENT
    avatar.line.fill.background()

    # Navn
    add_text_box(slide, x - 0.5, 5, 3.5, 0.6, name, 24, PRIMARY, True, PP_ALIGN.CENTER)
    add_text_box(slide, x - 0.5, 5.5, 3.5, 0.5, "Niels Brock", 18, GRAY, False, PP_ALIGN.CENTER)

# Beskrivelse
add_text_box(slide, 1.5, 6.3, 10.333, 0.8, "To unge iværksættere med passion for digital marketing - vokset op i den digitale verden", 20, GRAY, False, PP_ALIGN.CENTER)

# ============ SLIDE 3: PROBLEMET ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Problemet", 44, WHITE, True)

# Stat boxes
for i, (x, number, label) in enumerate([
    (1.5, "73%", "af små virksomheder\nkæmper med online\nmarkedsføring"),
    (5.5, "10K+", "kroner koster\ntraditionelle bureauer\nom måneden"),
    (9.5, "0", "gennemsigtighed\ni branchen")
]):
    # Box
    stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(3), Inches(3.5))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = OFF_WHITE
    stat_box.line.color.rgb = RGBColor(230, 235, 232)
    stat_box.adjustments[0] = 0.1

    # Nummer
    add_text_box(slide, x, 2.6, 3, 1, number, 48, ACCENT, True, PP_ALIGN.CENTER)

    # Label
    label_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.8), Inches(2.6), Inches(1.5))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 4: LØSNINGEN ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Vores løsning", 44, WHITE, True)

solutions = [
    ("Gennemsigtig marketing", "til fair priser"),
    ("No cure, no pay", "du betaler for resultater"),
    ("Moderne værktøjer", "og intelligent brug af AI"),
    ("Personlig service", "hvor du er prioriteten")
]

for i, (title, subtitle) in enumerate(solutions):
    y = 2.2 + (i * 1.2)

    # Ikon box
    icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(0.6), Inches(0.6))
    icon.fill.solid()
    icon.fill.fore_color.rgb = ACCENT
    icon.line.fill.background()
    icon.adjustments[0] = 0.25

    # Tekst
    add_text_box(slide, 2, y, 5, 0.6, title, 24, PRIMARY, True)
    add_text_box(slide, 2, y + 0.45, 5, 0.5, subtitle, 18, ACCENT, False)

# ============ SLIDE 5: SERVICES ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Hvad vi tilbyder", 44, WHITE, True)

services = [
    (0.8, 2.2, "Meta Ads", "Facebook & Instagram annoncer", META_BLUE),
    (4.05, 2.2, "Google Ads", "Søgemaskine annoncering", GOOGLE_BLUE),
    (7.3, 2.2, "SEO", "Søgemaskineoptimering", PRIMARY),
    (10.55, 2.2, "Hjemmesider", "Moderne webdesign", WEB_PURPLE)
]

for x, y, title, desc, color in services:
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.75), Inches(4))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(230, 235, 232)
    card.adjustments[0] = 0.08

    # Ikon
    icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.8), Inches(y + 0.5), Inches(1.15), Inches(1.15))
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()
    icon.adjustments[0] = 0.2

    # Titel
    add_text_box(slide, x, y + 2, 2.75, 0.6, title, 22, PRIMARY, True, PP_ALIGN.CENTER)

    # Beskrivelse
    desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.6), Inches(2.35), Inches(1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 6: HVORFOR OS ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Hvorfor vælge os?", 44, WHITE, True)

usps = [
    (0.8, 2.2, "No cure, no pay", "Du betaler kun for resultater"),
    (7, 2.2, "Komplet pakke", "Marketing OG hjemmeside"),
    (0.8, 4.5, "Digital natives", "Vi forstår online trends"),
    (7, 4.5, "Personlig service", "Du er vores førsteprioritet")
]

for x, y, title, desc in usps:
    # Card med venstre border
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.5), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(230, 235, 232)
    card.adjustments[0] = 0.06

    # Accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.15), Inches(0.12), Inches(1.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()

    # Ikon
    icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.4), Inches(y + 0.4), Inches(1), Inches(1))
    icon.fill.solid()
    icon.fill.fore_color.rgb = PRIMARY
    icon.line.fill.background()
    icon.adjustments[0] = 0.2

    # Tekst
    add_text_box(slide, x + 1.6, y + 0.4, 3.5, 0.5, title, 22, PRIMARY, True)
    add_text_box(slide, x + 1.6, y + 0.95, 3.5, 0.5, desc, 16, GRAY, False)

# ============ SLIDE 7: PROCESSEN ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY
header.line.fill.background()

add_text_box(slide, 0.8, 0.5, 10, 1, "Sådan arbejder vi", 44, WHITE, True)

# Forbindelseslinje
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.3), Inches(9.5), Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

steps = [
    (1.3, "1", "Gratis\nkonsultation"),
    (4.3, "2", "Strategi &\nforslag"),
    (7.3, "3", "Justering"),
    (10.3, "4", "Resultater")
]

for x, num, label in steps:
    # Cirkel
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.5), Inches(1.7), Inches(1.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()

    # Nummer
    add_text_box(slide, x, 2.85, 1.7, 1, num, 40, WHITE, True, PP_ALIGN.CENTER)

    # Label
    label_box = slide.shapes.add_textbox(Inches(x - 0.3), Inches(4.5), Inches(2.3), Inches(1.2))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 8: CTA ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, PRIMARY, ACCENT)

add_text_box(slide, 0.5, 2.2, 12.333, 1, "Klar til at vokse online?", 52, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 3.5, 12.333, 0.8, "Book en gratis og uforpligtende konsultation", 26, ACCENT_LIGHT, False, PP_ALIGN.CENTER)

# Kontakt boxes
contacts = [
    (3, "kontakt@nordicmarketing.dk"),
    (7.5, "nordicmarketing.dk")
]

for x, text in contacts:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5), Inches(3.5), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.fill.fore_color.brightness = 0.9
    box.line.fill.background()
    box.adjustments[0] = 0.3

    add_text_box(slide, x, 5.2, 3.5, 0.5, text, 16, PRIMARY, False, PP_ALIGN.CENTER)

# Gem
output_path = '/Users/theodorhauch/nordic-marketing/Nordic_Marketing_Pitch.pptx'
prs.save(output_path)
print(f"PowerPoint gemt: {output_path}")
