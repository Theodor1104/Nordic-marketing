from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Farver
PRIMARY = RGBColor(26, 61, 46)      # Mørkegrøn
ACCENT = RGBColor(61, 139, 110)     # Lysere grøn
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(74, 90, 81)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Baggrund
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Undertitel
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, bullets, has_icon=False):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Grøn header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(18)

    return slide

def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Venstre kolonne
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.5), Inches(0.6))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    # Højre kolonne
    right_header = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.5), Inches(0.6))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.5), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    return slide

def add_cta_slide(title, subtitle, contact_info):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Baggrund
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Undertitel
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Kontakt info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(1.5))
    tf = contact_box.text_frame
    for i, info in enumerate(contact_info):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = info
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    return slide

# === SLIDES ===

# Slide 1: Forside
add_title_slide(
    "Nordic Marketing",
    "Digital marketing der skaber resultater"
)

# Slide 2: Hvem er vi
add_content_slide(
    "Hvem er vi?",
    [
        "To studerende fra Niels Brock med passion for digital marketing",
        "Vokset op i den digitale verden med naturlig forståelse for online trends",
        "Praktisk erfaring gennem cases som Joe and the Juice",
        "Vi kombinerer friske idéer med professionel eksekvering"
    ]
)

# Slide 3: Vores mission
add_content_slide(
    "Vores mission",
    [
        "Gøre marketing mere gennemsigtigt og tilgængeligt",
        "Levere målbare resultater til fair priser",
        "Bryde med traditionelle, dyre bureauer",
        "Bruge moderne værktøjer og AI intelligent"
    ]
)

# Slide 4: Services
add_two_column_slide(
    "Hvad vi tilbyder",
    "Marketing",
    [
        "Meta Ads (Facebook & Instagram)",
        "Google Ads & Shopping",
        "SEO - Søgemaskineoptimering",
        "Content & strategi"
    ],
    "Web",
    [
        "Professionelle hjemmesider",
        "Responsivt design",
        "SEO-venlig struktur",
        "Analytics & tracking"
    ]
)

# Slide 5: Hvorfor vælge os
add_content_slide(
    "Hvorfor vælge Nordic Marketing?",
    [
        "No cure, no pay - du betaler kun for resultater",
        "Komplet løsning - marketing OG hjemmeside samlet",
        "Digital natives - vi forstår online trends og unge målgrupper",
        "Personlig service - du er vores prioritet",
        "Fair priser - professionel kvalitet uden bureaupriser"
    ]
)

# Slide 6: Vores proces
add_content_slide(
    "Sådan arbejder vi",
    [
        "1. Gratis konsultation - vi lytter til dine udfordringer",
        "2. Strategi - vi præsenterer vores bedste løsning",
        "3. Justering - vi tilpasser efter dine ønsker",
        "4. Eksekvering - vi implementerer og leverer resultater",
        "5. Opfølgning - løbende optimering og rapportering"
    ]
)

# Slide 7: CTA
add_cta_slide(
    "Klar til at vokse online?",
    "Book en gratis og uforpligtende konsultation i dag",
    [
        "kontakt@nordicmarketing.dk",
        "linkedin.com/company/nordic-marketing",
        "nordicmarketing.dk"
    ]
)

# Gem filen
prs.save('/Users/theodorhauch/nordic-marketing/Nordic_Marketing_Pitch.pptx')
print("PowerPoint gemt som: Nordic_Marketing_Pitch.pptx")
